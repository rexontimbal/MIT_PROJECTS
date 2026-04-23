"""
Generate Final Defense PowerPoint Presentation
AGNES Hotspot Detection System — MIT Capstone, SNSU
Author: Rexon L. Timbal
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os, math

# ── Colour Palette ──────────────────────────────────────────────
PNP_BLUE      = RGBColor(0x00, 0x30, 0x87)   # #003087
DARK_BLUE     = RGBColor(0x00, 0x1D, 0x52)   # darker header bg
ACCENT_BLUE   = RGBColor(0x00, 0x52, 0xCC)   # #0052CC
LIGHT_BLUE    = RGBColor(0xE8, 0xF0, 0xFE)   # light fill
SKY_BLUE      = RGBColor(0x42, 0x85, 0xF4)   # accent
ACCENT_RED    = RGBColor(0xDC, 0x14, 0x3C)   # #DC143C
ACCENT_GOLD   = RGBColor(0xF5, 0xA6, 0x23)   # gold accent
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
BLACK         = RGBColor(0x1A, 0x1A, 0x2E)
GRAY_700      = RGBColor(0x37, 0x41, 0x51)
GRAY_500      = RGBColor(0x6B, 0x72, 0x80)
GRAY_300      = RGBColor(0xD1, 0xD5, 0xDB)
GRAY_100      = RGBColor(0xF3, 0xF4, 0xF6)
GREEN         = RGBColor(0x05, 0x96, 0x69)
BG_OFF_WHITE  = RGBColor(0xFA, 0xFA, 0xFC)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

# ── Helpers ─────────────────────────────────────────────────────
def add_shape(slide, left, top, width, height, fill_color=None, line_color=None, line_w=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.shadow.inherit = False
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_w:
            shape.line.width = line_w
    else:
        shape.line.fill.background()
    return shape

def add_rounded_rect(slide, left, top, width, height, fill_color=None, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.shadow.inherit = False
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape

def add_circle(slide, left, top, size, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.shadow.inherit = False
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=14,
                 bold=False, color=BLACK, alignment=PP_ALIGN.LEFT,
                 font_name='Calibri', line_spacing=1.15, anchor=MSO_ANCHOR.TOP):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.text_frame.word_wrap = True
    txBox.text_frame.auto_size = False
    tf = txBox.text_frame
    tf.paragraphs[0].text = text
    tf.paragraphs[0].font.size = Pt(font_size)
    tf.paragraphs[0].font.bold = bold
    tf.paragraphs[0].font.color.rgb = color
    tf.paragraphs[0].font.name = font_name
    tf.paragraphs[0].alignment = alignment
    tf.paragraphs[0].space_after = Pt(0)
    tf.paragraphs[0].line_spacing = Pt(font_size * line_spacing)
    try:
        txBox.text_frame.paragraphs[0].font.language_id = None
    except:
        pass
    return txBox

def add_multiline(slide, left, top, width, height, lines, font_size=14,
                  color=BLACK, alignment=PP_ALIGN.LEFT, font_name='Calibri',
                  line_spacing=1.3, bold=False, anchor=MSO_ANCHOR.TOP):
    """lines = list of (text, bold_override, color_override, size_override) or just strings"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.text_frame.word_wrap = True
    tf = txBox.text_frame
    for i, item in enumerate(lines):
        if isinstance(item, str):
            txt, b, c, s = item, bold, color, font_size
        else:
            txt = item[0]
            b = item[1] if len(item) > 1 else bold
            c = item[2] if len(item) > 2 else color
            s = item[3] if len(item) > 3 else font_size
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = txt
        p.font.size = Pt(s)
        p.font.bold = b
        p.font.color.rgb = c
        p.font.name = font_name
        p.alignment = alignment
        p.space_after = Pt(2)
        p.line_spacing = Pt(s * line_spacing)
    return txBox

def add_icon_card(slide, left, top, width, height, icon_text, title, desc,
                  icon_color=ACCENT_BLUE, title_color=DARK_BLUE, desc_color=GRAY_700):
    """Card with colored circle icon, title, and description."""
    card = add_rounded_rect(slide, left, top, width, height, fill_color=WHITE, line_color=GRAY_300)
    # Icon circle
    circle_size = Inches(0.55)
    cx = left + Inches(0.25)
    cy = top + Inches(0.25)
    circ = add_circle(slide, cx, cy, circle_size, icon_color)
    add_text_box(slide, cx, cy + Inches(0.03), circle_size, circle_size,
                 icon_text, font_size=18, bold=True, color=WHITE,
                 alignment=PP_ALIGN.CENTER, font_name='Segoe UI Emoji')
    # Title
    add_text_box(slide, left + Inches(0.25), top + Inches(0.95), width - Inches(0.5), Inches(0.35),
                 title, font_size=13, bold=True, color=title_color)
    # Description
    add_text_box(slide, left + Inches(0.25), top + Inches(1.3), width - Inches(0.5), height - Inches(1.55),
                 desc, font_size=10, color=desc_color, line_spacing=1.25)
    return card

def slide_bg(slide, color=BG_OFF_WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_top_bar(slide, title_text, subtitle_text=None):
    """Dark blue top bar with title."""
    add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.15), fill_color=DARK_BLUE)
    # Gold accent line
    add_shape(slide, Inches(0), Inches(1.15), SLIDE_W, Inches(0.04), fill_color=ACCENT_GOLD)
    add_text_box(slide, Inches(0.7), Inches(0.2), Inches(10), Inches(0.5),
                 title_text, font_size=28, bold=True, color=WHITE, font_name='Calibri Light')
    if subtitle_text:
        add_text_box(slide, Inches(0.7), Inches(0.65), Inches(10), Inches(0.4),
                     subtitle_text, font_size=14, color=ACCENT_GOLD, font_name='Calibri')
    # Slide number placeholder area (right)
    return slide

def add_bullet_list(slide, left, top, width, height, items, font_size=13,
                    color=GRAY_700, bullet_color=ACCENT_BLUE, spacing=1.35):
    """Items = list of strings. Each gets a bullet."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.text_frame.word_wrap = True
    tf = txBox.text_frame
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        # Use a run for the bullet color
        run_bullet = p.add_run()
        run_bullet.text = "●  "
        run_bullet.font.size = Pt(font_size - 2)
        run_bullet.font.color.rgb = bullet_color
        run_bullet.font.name = 'Calibri'
        run = p.add_run()
        run.text = item
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.name = 'Calibri'
        p.space_after = Pt(4)
        p.line_spacing = Pt(font_size * spacing)
    return txBox

def add_number_step(slide, left, top, num, title, desc, num_color=ACCENT_BLUE):
    """Numbered step with large number, title, and desc."""
    # Number circle
    circ = add_circle(slide, left, top, Inches(0.5), num_color)
    add_text_box(slide, left, top + Inches(0.05), Inches(0.5), Inches(0.45),
                 str(num), font_size=18, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.65), top + Inches(0.02), Inches(2.4), Inches(0.3),
                 title, font_size=13, bold=True, color=DARK_BLUE)
    add_text_box(slide, left + Inches(0.65), top + Inches(0.32), Inches(2.4), Inches(0.4),
                 desc, font_size=10, color=GRAY_500, line_spacing=1.2)


# ═══════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE SLIDE
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
# Full dark blue background
add_shape(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill_color=DARK_BLUE)

# Subtle geometric decoration — diagonal accent
add_shape(slide, Inches(0), Inches(0), Inches(0.08), SLIDE_H, fill_color=ACCENT_GOLD)
# Bottom gold strip
add_shape(slide, Inches(0), Inches(7.1), SLIDE_W, Inches(0.04), fill_color=ACCENT_GOLD)

# "CAPSTONE DEFENSE" label
add_text_box(slide, Inches(1.2), Inches(0.9), Inches(11), Inches(0.4),
             "MASTER IN INFORMATION TECHNOLOGY  ·  CAPSTONE DEFENSE",
             font_size=13, bold=True, color=ACCENT_GOLD, font_name='Calibri',
             alignment=PP_ALIGN.CENTER)

# Main Title
add_multiline(slide, Inches(1.2), Inches(1.7), Inches(11), Inches(2.5), [
    ("AI-Based Hotspot Detection and", True, WHITE, 34),
    ("Reporting System for Road Accident Analysis:", True, WHITE, 34),
    ("A Decision Support System", True, SKY_BLUE, 34),
], alignment=PP_ALIGN.CENTER, line_spacing=1.25)

# Separator line
add_shape(slide, Inches(5.5), Inches(4.3), Inches(2.3), Inches(0.03), fill_color=ACCENT_GOLD)

# Author info
add_multiline(slide, Inches(1.2), Inches(4.7), Inches(11), Inches(1.5), [
    ("Presented by", False, GRAY_300, 12),
    ("REXON L. TIMBAL", True, WHITE, 22),
    ("", False, WHITE, 8),
    ("Faculty of the Graduate School", False, GRAY_300, 13),
    ("Surigao del Norte State University", False, ACCENT_GOLD, 14),
], alignment=PP_ALIGN.CENTER, line_spacing=1.3)

# Bottom footer
add_text_box(slide, Inches(1.2), Inches(6.7), Inches(11), Inches(0.35),
             "April 2026  ·  PNP Caraga Region  ·  AGNES Clustering Algorithm",
             font_size=11, color=GRAY_500, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# SLIDE 2 — OUTLINE / AGENDA
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(slide)
add_top_bar(slide, "PRESENTATION OUTLINE")

items = [
    ("01", "Introduction & Background"),
    ("02", "Problem Statement"),
    ("03", "Research Gap"),
    ("04", "Objectives"),
    ("05", "Methodology & Framework"),
    ("06", "System Architecture & Tech Stack"),
    ("07", "Key Features & Modules"),
    ("08", "AGNES Algorithm & Clustering"),
    ("09", "Results & Validation"),
    ("10", "Conclusion & Recommendations"),
]

col1_x = Inches(0.9)
col2_x = Inches(7.0)
y_start = Inches(1.7)
row_h = Inches(0.52)

for i, (num, label) in enumerate(items):
    col_x = col1_x if i < 5 else col2_x
    row_y = y_start + (i % 5) * row_h
    # Number
    add_text_box(slide, col_x, row_y, Inches(0.6), Inches(0.4),
                 num, font_size=20, bold=True, color=ACCENT_BLUE, font_name='Calibri Light')
    # Line
    add_shape(slide, col_x + Inches(0.65), row_y + Inches(0.22), Inches(0.4), Inches(0.02), fill_color=ACCENT_GOLD)
    # Label
    add_text_box(slide, col_x + Inches(1.15), row_y + Inches(0.03), Inches(4.5), Inches(0.4),
                 label, font_size=15, color=GRAY_700)


# ═══════════════════════════════════════════════════════════════
# SLIDE 3 — INTRODUCTION & BACKGROUND
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(slide)
add_top_bar(slide, "INTRODUCTION & BACKGROUND", "The Global and Local Context of Road Traffic Accidents")

# Stat cards row
stats = [
    ("1.19M", "Lives Lost Annually", "WHO Global Road\nSafety Report", ACCENT_RED),
    ("3%", "GDP Cost per Nation", "Economic impact of\nroad traffic injuries", ACCENT_BLUE),
    ("20–40", "Most Affected Age", "Working-age males\ndisproportionately affected", ACCENT_GOLD),
    ("Caraga", "Regional Focus", "Driver negligence as\nprimary cause of collisions", GREEN),
]

for i, (big_num, title, desc, accent) in enumerate(stats):
    cx = Inches(0.7) + i * Inches(3.15)
    cy = Inches(1.6)
    card = add_rounded_rect(slide, cx, cy, Inches(2.85), Inches(1.65), fill_color=WHITE, line_color=GRAY_300)
    # Top accent bar
    add_shape(slide, cx, cy, Inches(2.85), Inches(0.05), fill_color=accent)
    # Big number
    add_text_box(slide, cx + Inches(0.25), cy + Inches(0.25), Inches(2.4), Inches(0.55),
                 big_num, font_size=32, bold=True, color=accent)
    add_text_box(slide, cx + Inches(0.25), cy + Inches(0.8), Inches(2.4), Inches(0.3),
                 title, font_size=12, bold=True, color=DARK_BLUE)
    add_text_box(slide, cx + Inches(0.25), cy + Inches(1.1), Inches(2.4), Inches(0.5),
                 desc, font_size=10, color=GRAY_500, line_spacing=1.2)

# Bottom context paragraph
add_text_box(slide, Inches(0.7), Inches(3.6), Inches(12), Inches(0.7),
             "Road traffic accidents remain a critical public health and economic concern globally. "
             "In the Philippines, specifically the Caraga Region, high-risk neighborhoods such as Libertad, Butuan City "
             "experience frequent vehicular incidents, yet lack sophisticated analytical tools for systematic hotspot identification.",
             font_size=13, color=GRAY_700, line_spacing=1.4)

# Key insight box
insight_box = add_rounded_rect(slide, Inches(0.7), Inches(4.6), Inches(12), Inches(0.9),
                                fill_color=LIGHT_BLUE, line_color=ACCENT_BLUE)
add_text_box(slide, Inches(0.95), Inches(4.7), Inches(0.6), Inches(0.4),
             "💡", font_size=22, alignment=PP_ALIGN.CENTER, font_name='Segoe UI Emoji')
add_text_box(slide, Inches(1.5), Inches(4.75), Inches(10.8), Inches(0.7),
             "Current methods rely on manual reporting and generalized statistics, lacking the precision needed "
             "for targeted, data-driven traffic safety interventions.",
             font_size=13, bold=True, color=PNP_BLUE, line_spacing=1.35)


# ═══════════════════════════════════════════════════════════════
# SLIDE 4 — PROBLEM STATEMENT
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(slide)
add_top_bar(slide, "PROBLEM STATEMENT", "Identifying the Core Challenges")

problems = [
    ("📊", "Lack of Systematic Analysis",
     "Traffic accident data in the Caraga Region is collected manually without spatial clustering or pattern analysis, "
     "limiting the ability to identify recurring high-risk zones."),
    ("📍", "No Geospatial Hotspot Detection",
     "Existing processes do not leverage geographic coordinates to detect accident-prone areas. "
     "Interventions are reactive rather than proactive and preventive."),
    ("📋", "Fragmented Reporting Workflow",
     "Accident reports follow paper-based or unstructured digital formats (IRF), "
     "making data inconsistent and difficult to aggregate for analysis."),
    ("📉", "Absence of Decision Support Tools",
     "Law enforcement and traffic management agencies lack a centralized dashboard "
     "to visualize trends, severity, and prioritize resource allocation."),
]

for i, (icon, title, desc) in enumerate(problems):
    row = i // 2
    col = i % 2
    cx = Inches(0.7) + col * Inches(6.3)
    cy = Inches(1.6) + row * Inches(2.6)
    card = add_rounded_rect(slide, cx, cy, Inches(5.9), Inches(2.2), fill_color=WHITE, line_color=GRAY_300)
    # Icon
    add_text_box(slide, cx + Inches(0.3), cy + Inches(0.25), Inches(0.6), Inches(0.5),
                 icon, font_size=26, font_name='Segoe UI Emoji')
    # Title
    add_text_box(slide, cx + Inches(1.0), cy + Inches(0.3), Inches(4.5), Inches(0.35),
                 title, font_size=15, bold=True, color=DARK_BLUE)
    # Description
    add_text_box(slide, cx + Inches(1.0), cy + Inches(0.75), Inches(4.5), Inches(1.3),
                 desc, font_size=12, color=GRAY_700, line_spacing=1.35)


# ═══════════════════════════════════════════════════════════════
# SLIDE 5 — RESEARCH GAP
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(slide)
add_top_bar(slide, "RESEARCH GAP", "What Existing Studies Have Not Addressed")

# Left column — Current State
add_rounded_rect(slide, Inches(0.7), Inches(1.6), Inches(5.7), Inches(5.2),
                 fill_color=WHITE, line_color=GRAY_300)
add_shape(slide, Inches(0.7), Inches(1.6), Inches(5.7), Inches(0.06), fill_color=ACCENT_RED)
add_text_box(slide, Inches(1.0), Inches(1.85), Inches(5), Inches(0.4),
             "CURRENT STATE  ✕", font_size=16, bold=True, color=ACCENT_RED)

gap_current = [
    "Most studies focus on national-level or metro-area accident analysis, overlooking regional-level patterns in areas like Caraga.",
    "Existing systems use basic statistical summaries (counts, averages) without applying clustering algorithms for spatial pattern recognition.",
    "No integration of AGNES hierarchical clustering with a web-based decision support system for Philippine law enforcement.",
    "Paper-based Traffic Accident Reports (TAR) are difficult to digitize, aggregate, and analyze at scale.",
    "Lack of scientifically validated clustering results — no use of Silhouette, Davies-Bouldin, or Calinski-Harabasz metrics.",
]
for i, item in enumerate(gap_current):
    y = Inches(2.5) + i * Inches(0.82)
    add_text_box(slide, Inches(1.1), y, Inches(0.3), Inches(0.3),
                 "✕", font_size=13, bold=True, color=ACCENT_RED)
    add_text_box(slide, Inches(1.45), y, Inches(4.6), Inches(0.75),
                 item, font_size=11, color=GRAY_700, line_spacing=1.25)

# Right column — This Study
add_rounded_rect(slide, Inches(6.8), Inches(1.6), Inches(5.7), Inches(5.2),
                 fill_color=WHITE, line_color=GRAY_300)
add_shape(slide, Inches(6.8), Inches(1.6), Inches(5.7), Inches(0.06), fill_color=GREEN)
add_text_box(slide, Inches(7.1), Inches(1.85), Inches(5), Inches(0.4),
             "THIS STUDY ADDRESSES  ✓", font_size=16, bold=True, color=GREEN)

gap_this = [
    "Focuses specifically on the Caraga Region (5 provinces), capturing localized accident patterns and hotspots.",
    "Applies AGNES (Agglomerative Nesting) hierarchical clustering to geospatial accident coordinates for precise hotspot detection.",
    "Integrates clustering engine with a full-featured web-based DSS including maps, dashboards, and TAR-compliant reporting.",
    "Digitizes the entire TAR workflow with structured data capture, verification pipeline, and role-based access.",
    "Validates clustering quality using three scientific metrics: Silhouette Score, Davies-Bouldin Index, and Calinski-Harabasz Score.",
]
for i, item in enumerate(gap_this):
    y = Inches(2.5) + i * Inches(0.82)
    add_text_box(slide, Inches(7.2), y, Inches(0.3), Inches(0.3),
                 "✓", font_size=13, bold=True, color=GREEN)
    add_text_box(slide, Inches(7.55), y, Inches(4.6), Inches(0.75),
                 item, font_size=11, color=GRAY_700, line_spacing=1.25)


# ═══════════════════════════════════════════════════════════════
# SLIDE 6 — OBJECTIVES
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(slide)
add_top_bar(slide, "RESEARCH OBJECTIVES", "General and Specific Objectives")

# General Objective box
go_box = add_rounded_rect(slide, Inches(0.7), Inches(1.6), Inches(12), Inches(1.3),
                           fill_color=LIGHT_BLUE, line_color=ACCENT_BLUE)
add_text_box(slide, Inches(1.0), Inches(1.7), Inches(2), Inches(0.35),
             "GENERAL OBJECTIVE", font_size=11, bold=True, color=ACCENT_BLUE)
add_text_box(slide, Inches(1.0), Inches(2.05), Inches(11.4), Inches(0.7),
             "To develop an AI-based hotspot detection and reporting system that utilizes the AGNES clustering algorithm "
             "to identify high-risk road accident zones in the Caraga Region, providing a decision support tool for "
             "PNP traffic management and accident prevention.",
             font_size=13, color=DARK_BLUE, line_spacing=1.35)

# Specific objectives
specifics = [
    "Design and implement a web-based Traffic Accident Report (TAR) management system with structured data capture and multi-stage verification workflow.",
    "Apply the AGNES hierarchical clustering algorithm to geospatial accident data for identifying and ranking accident hotspots by severity.",
    "Develop an interactive analytics dashboard with real-time statistics, map visualizations, and trend analysis for data-driven decision making.",
    "Validate clustering results using scientific metrics (Silhouette Score, Davies-Bouldin Index, Calinski-Harabasz Score) to ensure reliability.",
    "Deploy the system as a responsive web platform with mobile application support for field accessibility by PNP personnel.",
]

add_text_box(slide, Inches(1.0), Inches(3.2), Inches(3), Inches(0.35),
             "SPECIFIC OBJECTIVES", font_size=11, bold=True, color=ACCENT_BLUE)

for i, obj in enumerate(specifics):
    y = Inches(3.7) + i * Inches(0.72)
    circ = add_circle(slide, Inches(1.0), y + Inches(0.02), Inches(0.35), ACCENT_BLUE)
    add_text_box(slide, Inches(1.0), y + Inches(0.02), Inches(0.35), Inches(0.35),
                 str(i+1), font_size=12, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1.5), y, Inches(11), Inches(0.65),
                 obj, font_size=12, color=GRAY_700, line_spacing=1.3)


# ═══════════════════════════════════════════════════════════════
# SLIDE 7 — METHODOLOGY & FRAMEWORK
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(slide)
add_top_bar(slide, "METHODOLOGY & FRAMEWORK", "Research Design and Development Approach")

# Methodology steps as a horizontal flow
steps = [
    ("1", "Data\nCollection", "Gathered CY2019-2024\ntraffic accident records\nfrom PNP Caraga"),
    ("2", "System\nDesign", "Designed TAR-based\nreporting workflow,\ndatabase schema, UI/UX"),
    ("3", "Development", "Built Django web app\nwith PostgreSQL, REST\nAPI, Leaflet.js maps"),
    ("4", "AGNES\nClustering", "Implemented hierarchical\nclustering on spatial\ncoordinates"),
    ("5", "Validation", "Applied Silhouette, D-B,\nC-H metrics to assess\nclustering quality"),
    ("6", "Testing &\nDeployment", "UAT with PNP users,\nIT expert evaluation,\nRailway deployment"),
]

for i, (num, title, desc) in enumerate(steps):
    cx = Inches(0.45) + i * Inches(2.12)
    cy = Inches(1.8)
    # Card
    card = add_rounded_rect(slide, cx, cy, Inches(1.95), Inches(3.0), fill_color=WHITE, line_color=GRAY_300)
    # Number circle
    circ = add_circle(slide, cx + Inches(0.65), cy + Inches(0.2), Inches(0.55), ACCENT_BLUE)
    add_text_box(slide, cx + Inches(0.65), cy + Inches(0.25), Inches(0.55), Inches(0.45),
                 num, font_size=20, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    # Title
    add_text_box(slide, cx + Inches(0.15), cy + Inches(0.9), Inches(1.65), Inches(0.65),
                 title, font_size=13, bold=True, color=DARK_BLUE, alignment=PP_ALIGN.CENTER,
                 line_spacing=1.15)
    # Desc
    add_text_box(slide, cx + Inches(0.15), cy + Inches(1.6), Inches(1.65), Inches(1.2),
                 desc, font_size=10, color=GRAY_500, alignment=PP_ALIGN.CENTER, line_spacing=1.25)
    # Arrow connector (except last)
    if i < len(steps) - 1:
        add_text_box(slide, cx + Inches(1.9), cy + Inches(1.2), Inches(0.3), Inches(0.4),
                     "▸", font_size=22, bold=True, color=ACCENT_GOLD, alignment=PP_ALIGN.CENTER)

# Bottom note
add_text_box(slide, Inches(0.7), Inches(5.2), Inches(12), Inches(0.5),
             "Research Design: Applied Research  |  Development Model: Agile / Iterative  |  "
             "Evaluation: User Acceptance Testing + IT Expert Review",
             font_size=12, color=GRAY_500, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# SLIDE 8 — SYSTEM ARCHITECTURE & TECH STACK
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(slide)
add_top_bar(slide, "SYSTEM ARCHITECTURE & TECHNOLOGY STACK")

# Architecture layers — horizontal blocks
layers = [
    ("PRESENTATION LAYER", "Responsive Web UI  ·  Leaflet.js Maps  ·  Chart.js Analytics  ·  Dark Mode  ·  Mobile App (Capacitor)", SKY_BLUE),
    ("APPLICATION LAYER", "Django 5.0.6  ·  REST API (DRF)  ·  Role-Based Access  ·  TAR Report Workflow  ·  Celery Task Queue", ACCENT_BLUE),
    ("INTELLIGENCE LAYER", "AGNES Clustering (SciPy)  ·  Severity Scoring  ·  Validation Metrics (scikit-learn)  ·  Geospatial Analysis", GREEN),
    ("DATA LAYER", "PostgreSQL 14  ·  Cloudinary Media  ·  CSV Import/Export  ·  Django ORM  ·  Redis Cache", PNP_BLUE),
]

for i, (label, tech, color) in enumerate(layers):
    cy = Inches(1.55) + i * Inches(1.15)
    # Colored label
    add_rounded_rect(slide, Inches(0.7), cy, Inches(2.6), Inches(0.85), fill_color=color)
    add_text_box(slide, Inches(0.85), cy + Inches(0.25), Inches(2.3), Inches(0.4),
                 label, font_size=11, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    # Tech detail
    add_rounded_rect(slide, Inches(3.5), cy, Inches(9.1), Inches(0.85), fill_color=WHITE, line_color=color)
    add_text_box(slide, Inches(3.8), cy + Inches(0.25), Inches(8.5), Inches(0.4),
                 tech, font_size=12, color=GRAY_700)

# Deployment box
add_rounded_rect(slide, Inches(0.7), Inches(6.05), Inches(11.9), Inches(0.75),
                 fill_color=LIGHT_BLUE, line_color=ACCENT_BLUE)
add_text_box(slide, Inches(1.0), Inches(6.15), Inches(2), Inches(0.35),
             "DEPLOYMENT", font_size=10, bold=True, color=ACCENT_BLUE)
add_text_box(slide, Inches(3.0), Inches(6.2), Inches(9), Inches(0.35),
             "Railway.app (Cloud)  ·  Gunicorn + Uvicorn  ·  WhiteNoise Static  ·  SSL/TLS  ·  Asia/Manila TZ",
             font_size=12, color=GRAY_700)


# ═══════════════════════════════════════════════════════════════
# SLIDE 9 — KEY FEATURES (1 of 2)
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(slide)
add_top_bar(slide, "KEY FEATURES & MODULES", "Core System Capabilities — Part 1")

features1 = [
    ("📍", "Hotspot Detection Engine", "AGNES hierarchical clustering on GPS coordinates with configurable linkage, threshold, and minimum cluster size.", ACCENT_RED),
    ("📊", "Analytics Dashboard", "Real-time stats (today/week/month), 12-month trends, province breakdown, time-of-day patterns, top 5 hotspots.", ACCENT_BLUE),
    ("🗺️", "Interactive Map", "Leaflet.js with marker clustering, fullscreen mode, satellite imagery, and hotspot zone visualization.", GREEN),
    ("📋", "TAR Reporting", "Structured Traffic Accident Report with 8-section form (Officer → Action Taken), multi-stage verification.", ACCENT_GOLD),
    ("👥", "Role-Based Access", "Officer, Admin, Supervisor, Data Manager roles with granular permissions and audit logging.", PNP_BLUE),
    ("📱", "Mobile Application", "Capacitor-based Android APK for field reporting by PNP personnel with offline-capable design.", SKY_BLUE),
]

for i, (icon, title, desc, color) in enumerate(features1):
    row = i // 3
    col = i % 3
    cx = Inches(0.55) + col * Inches(4.2)
    cy = Inches(1.6) + row * Inches(2.7)
    add_icon_card(slide, cx, cy, Inches(3.9), Inches(2.35), icon, title, desc, icon_color=color)


# ═══════════════════════════════════════════════════════════════
# SLIDE 10 — KEY FEATURES (2 of 2)
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(slide)
add_top_bar(slide, "KEY FEATURES & MODULES", "Core System Capabilities — Part 2")

features2 = [
    ("⚖️", "Severity Scoring", "Multi-factor scoring (0–100): frequency up to 40 pts, casualties up to 60 pts (10/fatality, 5/injury, 1/property damage).", ACCENT_RED),
    ("✅", "Clustering Validation", "Scientific quality assessment using Silhouette Score, Davies-Bouldin Index, and Calinski-Harabasz Score.", GREEN),
    ("📄", "Report Workflow", "Pending → Verified → Under Investigation → Resolved/Rejected pipeline with officer assignments and evidence uploads.", ACCENT_BLUE),
    ("📥", "Data Import/Export", "CSV bulk import (CY2019-2024 dataset), PDF/Excel export, data validation on ingest.", ACCENT_GOLD),
]

for i, (icon, title, desc, color) in enumerate(features2):
    col = i
    cx = Inches(0.45) + col * Inches(3.2)
    cy = Inches(1.6)
    add_icon_card(slide, cx, cy, Inches(2.95), Inches(2.6), icon, title, desc, icon_color=color)

# Workflow diagram representation
add_text_box(slide, Inches(0.7), Inches(4.8), Inches(12), Inches(0.4),
             "REPORT VERIFICATION WORKFLOW", font_size=12, bold=True, color=DARK_BLUE, alignment=PP_ALIGN.CENTER)

wf_steps = ["Submitted", "Pending\nVerification", "Verified", "Under\nInvestigation", "Resolved"]
wf_colors = [GRAY_500, ACCENT_GOLD, ACCENT_BLUE, PNP_BLUE, GREEN]
for i, (label, color) in enumerate(zip(wf_steps, wf_colors)):
    cx = Inches(1.2) + i * Inches(2.3)
    cy = Inches(5.3)
    add_rounded_rect(slide, cx, cy, Inches(1.8), Inches(0.85), fill_color=color)
    add_text_box(slide, cx, cy + Inches(0.12), Inches(1.8), Inches(0.65),
                 label, font_size=11, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER, line_spacing=1.15)
    if i < len(wf_steps) - 1:
        add_text_box(slide, cx + Inches(1.75), cy + Inches(0.2), Inches(0.6), Inches(0.4),
                     "→", font_size=22, bold=True, color=ACCENT_GOLD, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# SLIDE 11 — AGNES ALGORITHM
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(slide)
add_top_bar(slide, "AGNES CLUSTERING ALGORITHM", "Agglomerative Nesting — Hierarchical Clustering")

# Left: Algorithm steps
add_rounded_rect(slide, Inches(0.7), Inches(1.6), Inches(6.0), Inches(5.2),
                 fill_color=WHITE, line_color=GRAY_300)
add_text_box(slide, Inches(1.0), Inches(1.75), Inches(5), Inches(0.35),
             "HOW AGNES WORKS", font_size=14, bold=True, color=DARK_BLUE)

algo_steps = [
    ("1", "Initialize", "Each accident location (lat, lng) starts as its own individual cluster."),
    ("2", "Compute Distance", "Calculate pairwise distances between all clusters using Euclidean distance on coordinates."),
    ("3", "Merge Nearest", "Find the two closest clusters and merge them into one, using the selected linkage method (Complete/Single/Average)."),
    ("4", "Repeat", "Continue merging until the distance threshold is exceeded or all points form one cluster."),
    ("5", "Cut Dendrogram", "Apply distance threshold (≈0.05° ≈ 5km) to cut the hierarchy into flat clusters."),
    ("6", "Filter & Score", "Remove clusters below minimum size (3). Calculate severity scores based on casualties."),
]

for i, (num, title, desc) in enumerate(algo_steps):
    y = Inches(2.25) + i * Inches(0.72)
    circ = add_circle(slide, Inches(1.1), y, Inches(0.32), ACCENT_BLUE)
    add_text_box(slide, Inches(1.1), y + Inches(0.02), Inches(0.32), Inches(0.3),
                 num, font_size=11, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1.55), y, Inches(1.2), Inches(0.3),
                 title, font_size=11, bold=True, color=DARK_BLUE)
    add_text_box(slide, Inches(2.75), y, Inches(3.6), Inches(0.65),
                 desc, font_size=10, color=GRAY_700, line_spacing=1.2)

# Right: Configuration parameters
add_rounded_rect(slide, Inches(7.0), Inches(1.6), Inches(5.6), Inches(2.4),
                 fill_color=WHITE, line_color=GRAY_300)
add_shape(slide, Inches(7.0), Inches(1.6), Inches(5.6), Inches(0.05), fill_color=ACCENT_BLUE)
add_text_box(slide, Inches(7.3), Inches(1.8), Inches(5), Inches(0.35),
             "DEFAULT CONFIGURATION", font_size=13, bold=True, color=DARK_BLUE)

config_items = [
    ("Linkage Method:", "Complete (farthest neighbor)"),
    ("Distance Threshold:", "0.05° ≈ 5 km radius"),
    ("Min Cluster Size:", "3 accidents minimum"),
    ("Distance Metric:", "Euclidean"),
]
for i, (label, value) in enumerate(config_items):
    y = Inches(2.3) + i * Inches(0.4)
    add_text_box(slide, Inches(7.3), y, Inches(2.2), Inches(0.35),
                 label, font_size=11, bold=True, color=GRAY_700)
    add_text_box(slide, Inches(9.5), y, Inches(2.8), Inches(0.35),
                 value, font_size=11, color=ACCENT_BLUE)

# Severity scoring box
add_rounded_rect(slide, Inches(7.0), Inches(4.3), Inches(5.6), Inches(2.5),
                 fill_color=WHITE, line_color=GRAY_300)
add_shape(slide, Inches(7.0), Inches(4.3), Inches(5.6), Inches(0.05), fill_color=ACCENT_RED)
add_text_box(slide, Inches(7.3), Inches(4.5), Inches(5), Inches(0.35),
             "SEVERITY SCORING (0–100)", font_size=13, bold=True, color=DARK_BLUE)

severity_items = [
    ("Frequency Score:", "Up to 40 points", "Based on accident count"),
    ("Per Fatality:", "+10 points", "Maximum casualty weight"),
    ("Per Injury:", "+5 points", "Moderate casualty weight"),
    ("Property Damage:", "+1 point", "Minimal weight per incident"),
    ("Casualty Score:", "Up to 60 points", "Combined casualty total"),
]
for i, (label, value, note) in enumerate(severity_items):
    y = Inches(5.0) + i * Inches(0.33)
    add_text_box(slide, Inches(7.3), y, Inches(1.8), Inches(0.3),
                 label, font_size=10, bold=True, color=GRAY_700)
    add_text_box(slide, Inches(9.1), y, Inches(1.5), Inches(0.3),
                 value, font_size=10, bold=True, color=ACCENT_RED)
    add_text_box(slide, Inches(10.6), y, Inches(1.8), Inches(0.3),
                 note, font_size=9, color=GRAY_500)


# ═══════════════════════════════════════════════════════════════
# SLIDE 12 — RESULTS & VALIDATION
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(slide)
add_top_bar(slide, "RESULTS & VALIDATION", "Clustering Quality Metrics and Evaluation")

# Validation metrics explanation cards
metrics = [
    ("Silhouette Score", "−1 to +1", "Higher is better",
     "Measures how similar each point is to its own cluster versus neighboring clusters. "
     "Values close to +1 indicate well-separated, cohesive clusters.",
     ACCENT_BLUE),
    ("Davies-Bouldin Index", "0 to ∞", "Lower is better",
     "Measures average similarity between each cluster and its most similar one. "
     "Values close to 0 indicate clusters that are compact and well-separated.",
     GREEN),
    ("Calinski-Harabasz Score", "0 to ∞", "Higher is better",
     "Ratio of between-cluster dispersion to within-cluster dispersion. "
     "Higher values indicate denser, well-separated clustering.",
     ACCENT_GOLD),
]

for i, (name, rng, interpret, desc, color) in enumerate(metrics):
    cx = Inches(0.55) + i * Inches(4.2)
    cy = Inches(1.6)
    card = add_rounded_rect(slide, cx, cy, Inches(3.95), Inches(3.5), fill_color=WHITE, line_color=GRAY_300)
    add_shape(slide, cx, cy, Inches(3.95), Inches(0.06), fill_color=color)
    # Metric name
    add_text_box(slide, cx + Inches(0.25), cy + Inches(0.25), Inches(3.5), Inches(0.4),
                 name, font_size=16, bold=True, color=DARK_BLUE)
    # Range badge
    add_rounded_rect(slide, cx + Inches(0.25), cy + Inches(0.7), Inches(1.4), Inches(0.35), fill_color=LIGHT_BLUE)
    add_text_box(slide, cx + Inches(0.3), cy + Inches(0.72), Inches(1.3), Inches(0.3),
                 f"Range: {rng}", font_size=10, bold=True, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)
    # Interpretation badge
    add_rounded_rect(slide, cx + Inches(1.8), cy + Inches(0.7), Inches(1.6), Inches(0.35), fill_color=color)
    add_text_box(slide, cx + Inches(1.85), cy + Inches(0.72), Inches(1.5), Inches(0.3),
                 interpret, font_size=10, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    # Description
    add_text_box(slide, cx + Inches(0.25), cy + Inches(1.3), Inches(3.5), Inches(2.0),
                 desc, font_size=12, color=GRAY_700, line_spacing=1.4)

# Bottom evaluation summary
add_rounded_rect(slide, Inches(0.7), Inches(5.5), Inches(12), Inches(1.2),
                 fill_color=LIGHT_BLUE, line_color=ACCENT_BLUE)
add_text_box(slide, Inches(1.0), Inches(5.6), Inches(11.4), Inches(0.35),
             "EVALUATION METHODS", font_size=12, bold=True, color=ACCENT_BLUE)
add_text_box(slide, Inches(1.0), Inches(5.95), Inches(11.4), Inches(0.6),
             "User Acceptance Testing (UAT) with PNP Caraga personnel  ·  "
             "IT Expert Evaluation by domain specialists  ·  "
             "Clustering validation metrics applied to CY2019-2024 dataset across 5 Caraga provinces",
             font_size=12, color=GRAY_700, line_spacing=1.35)


# ═══════════════════════════════════════════════════════════════
# SLIDE 13 — SYSTEM DEMO / SCREENSHOTS (placeholder)
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(slide)
add_top_bar(slide, "SYSTEM DEMONSTRATION", "Live Demo / Screenshots")

# Placeholder boxes for screenshots
placeholders = [
    ("Dashboard & Analytics", Inches(0.7), Inches(1.6), Inches(5.8), Inches(3.5)),
    ("Hotspot Map", Inches(6.8), Inches(1.6), Inches(5.8), Inches(3.5)),
    ("TAR Report Form", Inches(0.7), Inches(5.3), Inches(3.8), Inches(1.7)),
    ("Clustering Results", Inches(4.7), Inches(5.3), Inches(3.8), Inches(1.7)),
    ("Mobile App", Inches(8.7), Inches(5.3), Inches(3.8), Inches(1.7)),
]

for label, x, y, w, h in placeholders:
    card = add_rounded_rect(slide, x, y, w, h, fill_color=WHITE, line_color=GRAY_300)
    # Center label
    add_text_box(slide, x, y + h/2 - Inches(0.35), w, Inches(0.35),
                 "[ Insert Screenshot ]", font_size=13, color=GRAY_500, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x, y + h/2, w, Inches(0.35),
                 label, font_size=14, bold=True, color=DARK_BLUE, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# SLIDE 14 — CONCLUSION
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(slide)
add_top_bar(slide, "CONCLUSION", "Summary of Contributions")

conclusions = [
    ("✓", "Successfully developed an AI-based decision support system that applies AGNES hierarchical clustering to identify road accident hotspots in the Caraga Region."),
    ("✓", "The system digitizes the entire Traffic Accident Report (TAR) workflow, replacing paper-based processes with a structured, verifiable digital pipeline."),
    ("✓", "Scientific validation through Silhouette Score, Davies-Bouldin Index, and Calinski-Harabasz Score confirms the reliability and quality of clustering results."),
    ("✓", "The interactive dashboard and map visualization provide law enforcement with actionable, data-driven insights for resource allocation and accident prevention."),
    ("✓", "Successful deployment as a responsive web platform with mobile application support demonstrates practical applicability for PNP Caraga field operations."),
]

for i, (check, text) in enumerate(conclusions):
    y = Inches(1.7) + i * Inches(0.95)
    card = add_rounded_rect(slide, Inches(0.7), y, Inches(12), Inches(0.8),
                             fill_color=WHITE, line_color=GRAY_300)
    # Green check circle
    circ = add_circle(slide, Inches(1.0), y + Inches(0.15), Inches(0.42), GREEN)
    add_text_box(slide, Inches(1.0), y + Inches(0.15), Inches(0.42), Inches(0.42),
                 check, font_size=16, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1.65), y + Inches(0.15), Inches(10.8), Inches(0.55),
                 text, font_size=12, color=GRAY_700, line_spacing=1.3)


# ═══════════════════════════════════════════════════════════════
# SLIDE 15 — RECOMMENDATIONS
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(slide)
add_top_bar(slide, "RECOMMENDATIONS", "Future Enhancements and Directions")

recs = [
    ("🔮", "Predictive Analytics",
     "Integrate time-series forecasting models to predict future accident hotspots based on historical trends and seasonal patterns."),
    ("🤖", "Advanced AI Models",
     "Explore deep learning approaches (e.g., DBSCAN, OPTICS) for comparison with AGNES, and integrate natural language processing for automated report analysis."),
    ("🌐", "Regional Expansion",
     "Extend the system to cover additional Philippine regions beyond Caraga, creating a nationwide accident analysis platform."),
    ("📡", "Real-Time Data Integration",
     "Connect with CCTV feeds, IoT sensors, and traffic management systems for real-time accident detection and automatic reporting."),
    ("🔗", "Inter-Agency Integration",
     "Enable data sharing with DPWH, DOTr, LTO, and hospital systems for comprehensive road safety ecosystem analysis."),
]

for i, (icon, title, desc) in enumerate(recs):
    cx = Inches(0.7)
    cy = Inches(1.6) + i * Inches(1.1)
    # Icon
    add_text_box(slide, cx, cy + Inches(0.1), Inches(0.5), Inches(0.4),
                 icon, font_size=20, font_name='Segoe UI Emoji')
    # Title
    add_text_box(slide, cx + Inches(0.65), cy + Inches(0.05), Inches(3), Inches(0.35),
                 title, font_size=14, bold=True, color=DARK_BLUE)
    # Desc
    add_text_box(slide, cx + Inches(0.65), cy + Inches(0.4), Inches(11.5), Inches(0.6),
                 desc, font_size=12, color=GRAY_700, line_spacing=1.3)
    # Separator
    if i < len(recs) - 1:
        add_shape(slide, Inches(1.35), cy + Inches(0.95), Inches(11), Inches(0.01), fill_color=GRAY_300)


# ═══════════════════════════════════════════════════════════════
# SLIDE 16 — THANK YOU / Q&A
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
# Full dark blue bg
add_shape(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill_color=DARK_BLUE)
add_shape(slide, Inches(0), Inches(0), Inches(0.08), SLIDE_H, fill_color=ACCENT_GOLD)
add_shape(slide, Inches(0), Inches(7.1), SLIDE_W, Inches(0.04), fill_color=ACCENT_GOLD)

add_multiline(slide, Inches(1.2), Inches(1.5), Inches(11), Inches(2.0), [
    ("Thank You", True, WHITE, 48),
    ("", False, WHITE, 10),
    ("Questions & Discussion", False, SKY_BLUE, 24),
], alignment=PP_ALIGN.CENTER, line_spacing=1.3)

# Separator
add_shape(slide, Inches(5.5), Inches(3.8), Inches(2.3), Inches(0.03), fill_color=ACCENT_GOLD)

add_multiline(slide, Inches(1.2), Inches(4.2), Inches(11), Inches(2.5), [
    ("AI-Based Hotspot Detection and Reporting System", True, WHITE, 16),
    ("for Road Accident Analysis: A Decision Support System", False, GRAY_300, 14),
    ("", False, WHITE, 10),
    ("Rexon L. Timbal", True, ACCENT_GOLD, 18),
    ("Master in Information Technology", False, GRAY_300, 13),
    ("Surigao del Norte State University", False, WHITE, 13),
], alignment=PP_ALIGN.CENTER, line_spacing=1.3)


# ═══════════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════════
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                           "AGNES_Final_Defense_Presentation.pptx")
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
